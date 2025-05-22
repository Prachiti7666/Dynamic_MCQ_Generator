from flask import Flask, render_template, request, redirect, url_for, session, make_response,flash, Response,abort,send_from_directory,send_file,jsonify
from azure.core.credentials import AzureKeyCredential
from azure.storage.blob import BlobServiceClient
from azure.search.documents import SearchClient
from azure.search.documents.indexes import SearchIndexClient
import openai
import os
import uuid
import tempfile
import PyPDF2
from docx import Document
from pptx import Presentation
import pyodbc
import requests
from requests_toolbelt.multipart.encoder import MultipartEncoder
from bs4 import BeautifulSoup
from datetime import datetime
import pytz
from moviepy.editor import VideoFileClip
import math
import soundfile as sf
import azure.cognitiveservices.speech as speechsdk
from typing import Optional,List, Tuple
import pandas as pd
import string
import base64
from io import BytesIO
import numpy as np
from urllib.parse import urlparse, unquote
import re




# Initialize Flask app
app = Flask(__name__,static_folder='static')

# Load environment variables

AZURE_SQL_SERVER="pocgenai.database.windows.net"
AZURE_SQL_DATABASE="qnadb"
AZURE_SQL_USERNAME="AZURE_SQL_USERNAME"
AZURE_SQL_PASSWORD="AZURE_SQL_PASSWORD"

AZURE_SPEECH_KEY = "AZURE_SPEECH_KEY"

FLASK_SECRET_KEY = "FLASK_SECRET_KEY"

LOCAL_SQL_SERVER="LOCAL_SQL_SERVER"
LOCAL_SQL_DATABASE="LOCAL_SQL_DATABASE"

# App configuration
app.secret_key = FLASK_SECRET_KEY



# # Database connection string
# connection_string = (
#     f"Driver={{ODBC Driver 17 for SQL Server}};"
#     f"Server=tcp:{AZURE_SQL_SERVER},1433;"
#     f"Database={AZURE_SQL_DATABASE};"
#     f"Uid={AZURE_SQL_USERNAME};"
#     f"Pwd={AZURE_SQL_PASSWORD};"
#     "Encrypt=yes;TrustServerCertificate=no;"
#     "Connection Timeout=50;"  # Add timeout parameter
#     "Command Timeout=50;" 
#     "Pooling=true;"
#     "Max Pool Size=200;"
#     "Min Pool Size=10;"    # Add command timeout
# )

connection_string = (
    f"Driver={{ODBC Driver 17 for SQL Server}};"
    f"Server={LOCAL_SQL_SERVER};"
    f"Database={LOCAL_SQL_DATABASE};"
    "Trusted_Connection=yes;"
)

# Azure service configurations
AZURE_STORAGE_CONNECTION_STRING = "AZURE_STORAGE_CONNECTION_STRING"
AZURE_SEARCH_SERVICE_ENDPOINT = "AZURE_SEARCH_SERVICE_ENDPOINT"
AZURE_SEARCH_KEY = "AZURE_SEARCH_KEY"
AZURE_SEARCH_INDEX_NAME = "AZURE_SEARCH_INDEX_NAME"


SPEECH_SUBSCRIPTION_KEY = "SPEECH_SUBSCRIPTION_KEY"
SPEECH_REGION = "SPEECH_REGION"
SPEECH_ENDPOINT = "SPEECH_ENDPOINT"


# Initialize Azure clients
blob_service_client = BlobServiceClient.from_connection_string(AZURE_STORAGE_CONNECTION_STRING)
credential = AzureKeyCredential(AZURE_SEARCH_KEY)
search_client = SearchClient(
    endpoint=AZURE_SEARCH_SERVICE_ENDPOINT,
    index_name=AZURE_SEARCH_INDEX_NAME,
    credential=credential
)
search_index_client = SearchIndexClient(
    endpoint=AZURE_SEARCH_SERVICE_ENDPOINT,
    credential=credential
)

def get_db_connection():
    try:
        return pyodbc.connect(connection_string)
    except pyodbc.Error as ex:
        sqlstate = ex.args[1]
        print(f"Connection failed: {sqlstate}")
        return None

class User:
    @staticmethod
    def query_by_email(email):
        conn = get_db_connection()
        if conn is None:
            print("Database connection failed.")
            return None
        
        # Use a connection pool instead of creating new connections
        try:
            cursor = conn.cursor()
            # Only select the columns you actually need
            cursor.execute("SELECT id, password, role FROM Users WHERE email = ?", (email,))
            user = cursor.fetchone()
            return user
        except Exception as e:
            print(f"Database error: {e}")
            return None
        finally:
            cursor.close()
            conn.close()

    @staticmethod
    def query_by_id(user_id):
        conn = get_db_connection()
        cursor = conn.cursor()
        cursor.execute("SELECT * FROM Users WHERE id = ?", (user_id,))
        user = cursor.fetchone()
        cursor.close()
        conn.close()
        return user

# Utility functions
def extract_text_from_web(url):
    """Extract text content from a web page"""
    response = requests.get(url)
    soup = BeautifulSoup(response.content, 'html.parser')
    text = soup.get_text(separator=' ')
    return text

def extract_text_from_pdf(pdf_file):
    """Extract text content from PDF file"""
    pdf_reader = PyPDF2.PdfReader(pdf_file)
    text = ""
    for page in pdf_reader.pages:
        text += page.extract_text()
    return text

def extract_text_from_docx(file_path):
    doc = Document(file_path)
    text = []
    for paragraph in doc.paragraphs:
        text.append(paragraph.text)
    return '\n'.join(text)

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return '\n'.join(text)

def chunk_text(text, chunk_size=5000, overlap=100):
    """Split text into overlapping chunks."""
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        if end < len(text):
            last_period = text[start:end].rfind('.')
            last_newline = text[start:end].rfind('\n')
            break_point = max(last_period, last_newline)
            if break_point != -1:
                end = start + break_point + 1
        chunk = text[start:end]
        chunks.append(chunk.strip())
        start = end - overlap
    return chunks

def convert_utc_to_ist(utc_time):
    ist = pytz.timezone('Asia/Kolkata')
    utc_time = utc_time.replace(tzinfo=pytz.utc)
    ist_time = utc_time.astimezone(ist)
    return ist_time



def transcribe_audio(file_path, api_key):
    url = 'url'
    # Create a MultipartEncoder object
    m = MultipartEncoder(
        fields={
            'audio': ('YourAudioFile', open(file_path, 'rb'), 'audio/wav'),
            'definition': '{"locales":["en-US"]}'
        }
    )
    
    headers = {
        'Ocp-Apim-Subscription-Key': api_key,
        'Content-Type': m.content_type
    }

    response = requests.post(url, headers=headers, data=m)
    text_content = ""

    if response.status_code == 200:
        print("Transcription successful!")
        data = response.json()
        combined_phrases = data.get("combinedPhrases", [])
        if combined_phrases:
            print(combined_phrases[0]["text"])
            text_content = combined_phrases[0]["text"]
        else:
            print("No transcription found.")
    else:
        print(f"Failed to transcribe: {response.status_code} - {response.text}")

    return text_content


def generate_questions(context, num_questions, test_id):
    """Generate questions dynamically using Azure OpenAI and save them to the database"""
    openai.api_type = "azure"
    openai.api_base = "api_base"
    openai.api_version = "2024-05-01-preview"
    openai.api_key = "api_key"


    prompt = f"""Based on the following content, generate only {num_questions} multiple-choice questions. Don't generate more questions.
    Format each question with 4 options and indicate the correct answer.
    Content: {context}
    Generate questions from all parts of the content.
    Generate questions in the following format strictly:
    Q1: [Question text]
    A) [Option 1]
    B) [Option 2]
    C) [Option 3]
    D) [Option 4]
    Correct: [A/B/C/D]
    """

    response = openai.ChatCompletion.create(
        engine="gpt-4o",
        messages=[
            {"role": "system", "content": "You are a helpful assistant that generates multiple-choice questions."},
            {"role": "user", "content": prompt}
        ],
        temperature=0.7,
        max_tokens=6500
    )

    questions = []
    current_question = {}
    conn = get_db_connection()
    cursor = conn.cursor()

    for line in response.choices[0].message.content.split('\n'):
        if line.startswith('Q'):
            if current_question:
                questions.append(current_question)
            current_question = {
                'text': line.split(': ')[1],
                'options': [],
                'correct_answer': None
            }
        elif line.startswith(('A)', 'B)', 'C)', 'D)')):
            option = {'id': line[0], 'text': line[3:].strip()}
            current_question['options'].append(option)
        elif line.startswith('Correct:'):
            current_question['correct_answer'] = line.split(': ')[1].strip()

    if current_question:
        questions.append(current_question)

    for question in questions:
        cursor.execute("""
            INSERT INTO TestQuestions (
                id, test_id, question_text, option_a, option_b, option_c, option_d, correct_answer
            ) VALUES (
                NEWID(), ?, ?, ?, ?, ?, ?, ?
            )
        """, (
            test_id,
            question['text'],
            question['options'][0]['text'],
            question['options'][1]['text'],
            question['options'][2]['text'],
            question['options'][3]['text'],
            question['correct_answer']
        ))

    conn.commit()
    cursor.close()
    conn.close()
    return questions


app.config['UPLOAD_FOLDER'] = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'uploads')

# Create uploads directory if it doesn't exist
if not os.path.exists(app.config['UPLOAD_FOLDER']):
    os.makedirs(app.config['UPLOAD_FOLDER'])


# Routes
@app.route('/')
def home():
    return render_template('index.html')


@app.route('/bulk-register', methods=['POST'])
def bulk_register():
    if 'trainee_list' not in request.files or request.files['trainee_list'].filename == '':
        return jsonify({'error': 'Please upload a document'}), 400

    ALLOWED_EXCEL_EXTENSIONS = {'.xlsx', '.xls'}
    trainee_file = request.files['trainee_list']
    file_ext = os.path.splitext(trainee_file.filename)[1].lower()

    if file_ext not in ALLOWED_EXCEL_EXTENSIONS:
        return jsonify({'error': 'Invalid file format. Please use Excel (.xlsx or .xls) file.'}), 400

    try:
        # Read Excel file directly from memory
        df = pd.read_excel(trainee_file, engine='openpyxl')
        
        # Validate columns
        required_columns = {'Name', 'Email'}
        if not required_columns.issubset(df.columns):
            return jsonify({'error': 'The file must contain columns named "Name" and "Email"'}), 400

        # Clean data
        df = df.dropna(subset=['Email', 'Name'])
        if df.empty:
            return jsonify({'error': 'No valid data found in the file'}), 400

        # Get list of emails from the Excel file
        new_emails = df['Email'].tolist()

        # Configure connection for Azure SQL
        conn = get_db_connection()
        
        try:
            cursor = conn.cursor()
            
            # Fixed email check query using proper SQL Server syntax
            check_query = """
            SELECT email FROM Users 
            WHERE email IN ({})
            """.format(','.join("'" + email.replace("'", "''") + "'" for email in new_emails))
            
            cursor.execute(check_query)
            existing_emails = [row[0] for row in cursor.fetchall()]

            if existing_emails:
                # Format the list of existing emails
                email_list = ', '.join(existing_emails)
                return jsonify({'error': f'The following emails are already registered: {email_list}'}), 400

            # Generate passwords and UUIDs efficiently
            num_records = len(df)
            # df['Password'] = [''.join(np.random.choice(
            #     list(string.ascii_letters + string.digits), size=10)) 
            #     for _ in range(num_records)]
            df['Password'] = [''.join("Welcome@1") 
                for _ in range(num_records)]
            df['id'] = [str(uuid.uuid4()) for _ in range(num_records)]

            cursor.fast_executemany = True

            # Prepare the insert statement with proper parameter placeholders for SQL Server
            insert_query = """
            INSERT INTO Users (id, name, email, password, role)
            VALUES (?, ?, ?, ?, ?)
            """

            # Batch size optimization for Azure SQL
            BATCH_SIZE = 1000
            total_records = len(df)
            
            for i in range(0, total_records, BATCH_SIZE):
                batch_df = df.iloc[i:min(i + BATCH_SIZE, total_records)]
                
                # Prepare batch values
                batch_values = [(row['id'], row['Name'], row['Email'], 
                               row['Password'], 'trainee') 
                              for _, row in batch_df.iterrows()]
                
                # Execute batch insert
                cursor.executemany(insert_query, batch_values)
                conn.commit()

            # Prepare download file
            download_df = df[['Name', 'Email', 'Password']]
            
            # Generate Excel file in memory
            output = BytesIO()
            with pd.ExcelWriter(output, engine='openpyxl') as writer:
                download_df.to_excel(writer, index=False)
            output.seek(0)

            return send_file(
                output,
                mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                as_attachment=True,
                download_name=f'trainee_list_{convert_utc_to_ist(datetime.now()).strftime("%Y%m%d_%H%M%S")}.xlsx'
            )

        except Exception as e:
            conn.rollback()
            app.logger.error(f'Database error during bulk registration: {str(e)}')
            return jsonify({'error': 'Database error. Please try again.'}), 500
        finally:
            conn.close()

    except pd.errors.EmptyDataError:
        return jsonify({'error': 'The uploaded file is empty'}), 400
    except Exception as e:
        app.logger.error(f'Error processing bulk registration: {str(e)}')
        return jsonify({'error': 'Error processing the file. Please check the format and try again.'}), 500
    

    
@app.route('/register', methods=['GET', 'POST'])
def register():
    if request.method == 'POST':
        email = request.form.get('email')
        password = request.form.get('password')
        name = request.form.get('name')
        role = request.form.get('role')

        conn = get_db_connection()
        cursor = conn.cursor()

        cursor.execute("SELECT id FROM Users WHERE email = ?", (email,))
        if cursor.fetchone():
            cursor.close()
            conn.close()
            flash('Email address already registered','error_registration')
            return redirect(url_for('login'))

        cursor.execute("""
            INSERT INTO Users (id, email, password, name, role)
            VALUES (NEWID(), ?, ?, ?, ?)
        """, (email, password, name, role))

        conn.commit()
        cursor.close()
        conn.close()

        flash('Registration successful, please log in.','success')
        return redirect(url_for('login'))

    return render_template('register.html')

from contextlib import contextmanager

CONNECTION_POOL = None

def init_connection_pool(pool_size=5):
    global CONNECTION_POOL
    # Use a proper connection pooling library like SQLAlchemy
    # This is a simplified example
    CONNECTION_POOL = [create_connection() for _ in range(pool_size)]

def create_connection():
    # Your connection creation logic
    pass

@contextmanager
def get_connection():
    conn = CONNECTION_POOL.pop() if CONNECTION_POOL else create_connection()
    try:
        yield conn
    finally:
        if CONNECTION_POOL is not None:
            CONNECTION_POOL.append(conn)
        else:
            conn.close()

@app.route('/login', methods=['GET', 'POST'])
def login():
    if request.method == 'POST':
        email = request.form.get('email')
        password = request.form.get('password')
        
        # Add input validation here
        if not email or not password:
            flash('Email and password are required', 'error_login')
            return render_template('login.html')
            
        user = User.query_by_email(email)
        
        # Use a constant-time comparison for passwords to prevent timing attacks
        # Also, you should be using password hashing, not storing plaintext passwords
        if user and user.password == password:
            # Store essential user data in session to avoid future queries
            session['user_id'] = str(user.id)
            session['user_role'] = user.role
            
            # Simplified redirect logic
            dashboard = 'trainer_dashboard' if user.role == 'trainer' else 'trainee_dashboard'
            return redirect(url_for(dashboard))
        
        flash('Invalid email or password', 'error_login')
    return render_template('login.html')

@app.route('/logout')
def logout():
    session.clear()
    return redirect(url_for('home'))

@app.route('/trainer/dashboard')
def trainer_dashboard():
    if 'user_id' not in session:
        return redirect(url_for('login'))

    user = User.query_by_id(session['user_id'])
    if user is None or user.role != 'trainer':
        return redirect(url_for('login'))

    conn = get_db_connection()
    cursor = conn.cursor()

    cursor.execute("""
        SELECT 
            ts.id,
            ts.title,
            ts.created_at,
            ts.qualification_marks,
            ts.num_questions,
            COUNT(DISTINCT ta.id) as total_trainees,
            COUNT(DISTINCT r.id) as total_attempts,
            COUNT(DISTINCT CASE WHEN r.id IS NOT NULL THEN ta.trainee_id END) as attempted_trainees,
            SUM(CASE WHEN r.score >= ts.qualification_marks THEN 1 ELSE 0 END) as passed_count
        FROM TrainingSession ts
        LEFT JOIN TraineeAssignment ta ON ts.id = ta.session_id
        LEFT JOIN Results r ON ta.id = r.assignment_id
        WHERE ts.created_by = ?
        GROUP BY ts.id, ts.created_at, 
                 ts.qualification_marks, ts.num_questions, ts.title
        ORDER BY ts.created_at DESC
    """, (session['user_id'],))

    training_sessions = []
    columns = [column[0] for column in cursor.description]
    for row in cursor.fetchall():
        session_dict = dict(zip(columns, row))

        # Convert the created_at time from UTC to IST
        session_dict['created_at'] = convert_utc_to_ist(session_dict['created_at'])

        session_dict['trainees'] = []
        training_sessions.append(session_dict)

    for session_dict in training_sessions:
        cursor.execute("""
            SELECT
                u.name as trainee_name,
                u.email as trainee_email,
                r.score as latest_score,
                COUNT(r2.id) as attempt_count,
                CASE WHEN r.score >= ts.qualification_marks THEN 1 ELSE 0 END as has_passed,
                r.completed_at as last_attempt,
                ta.id as assignment_id,
                r.id as result_id
            FROM TraineeAssignment ta
            JOIN Users u ON ta.trainee_id = u.id
            JOIN TrainingSession ts ON ta.session_id = ts.id
            LEFT JOIN Results r ON ta.id = r.assignment_id AND r.completed_at = (
                SELECT MAX(completed_at) 
                FROM Results 
                WHERE assignment_id = ta.id
            )
            LEFT JOIN Results r2 ON ta.id = r2.assignment_id
            WHERE ta.session_id = ?
            GROUP BY u.name, u.email, ta.id, r.id, r.score, r.completed_at, ts.qualification_marks
        """, (session_dict['id'],))
        trainee_columns = [column[0] for column in cursor.description]
        trainee_rows = cursor.fetchall()
        for trainee_row in trainee_rows:
            trainee_dict = dict(zip(trainee_columns, trainee_row))
            # Convert the last_attempt time from UTC to IST if it exists
            if trainee_dict['last_attempt']:
                trainee_dict['last_attempt'] = convert_utc_to_ist(trainee_dict['last_attempt'])
            session_dict['trainees'].append(trainee_dict)

    cursor.close()
    conn.close()

    return render_template('trainer_dashboard.html', training_sessions=training_sessions)


def get_filename_or_title(url):
    parsed_url = urlparse(url)
    
    # Check if the URL points to a file
    if re.search(r'\.\w+$', parsed_url.path):
        filename = unquote(parsed_url.path.split('/')[-1])
        # Remove date and time prefix using regex
        filename = re.sub(r'^\d{8}_\d{6}_', '', filename)
        return filename
    else:
        # Fetch the web page and extract the title
        try:
            response = requests.get(url)
            response.raise_for_status()  # Check for HTTP errors
            soup = BeautifulSoup(response.content, 'html.parser')
            title = soup.title.string.strip()
            return title
        except requests.RequestException as e:
            return f"Error fetching page title: {e}"

@app.route('/trainee/dashboard')
def trainee_dashboard():
    if 'user_id' not in session:
        return redirect(url_for('login'))

    conn = get_db_connection()
    cursor = conn.cursor()



    # First, get all the sessions with their basic info
    cursor.execute("""
        SELECT 
            ts.id as session_id,
            ts.title,
            ts.qualification_marks,
            ts.num_questions,
            ts.created_at as test_created_at,
            ta.id as assignment_id,
            (SELECT COUNT(*) FROM Results WHERE assignment_id = ta.id) as attempt_count,
            (SELECT MAX(score) FROM Results WHERE assignment_id = ta.id) as best_score,
            STRING_AGG(td.document_url, ',') as document_urls,
            STRING_AGG(td.document_id, ',') as document_ids
        FROM TrainingSession ts
        JOIN TraineeAssignment ta ON ts.id = ta.session_id
        LEFT JOIN Training_Docs td ON ts.id = td.session_id
        WHERE ta.trainee_id = ?
        GROUP BY 
            ts.id, ts.title, ts.qualification_marks, ts.num_questions,
            ts.created_at, ta.id
        ORDER BY ts.created_at DESC
    """, (session['user_id'],))

    available_sessions = []
    columns = [column[0] for column in cursor.description]
    for row in cursor.fetchall():
        session_dict = dict(zip(columns, row))
        
        # Convert the test_created_at time from UTC to IST
        session_dict['test_created_at'] = convert_utc_to_ist(session_dict['test_created_at'])
        
        # Split the concatenated document URLs and IDs into lists
        session_dict['documents'] = []
        if session_dict.get('document_urls') and session_dict.get('document_ids'):
            urls = session_dict['document_urls'].split(',')
            filenames = [get_filename_or_title(url) for url in urls]
            ids = session_dict['document_ids'].split(',')
            session_dict['documents'] = list(zip(urls, filenames, ids))
        
        # Get all attempts for this session
        cursor.execute("""
            SELECT 
                r.id,
                r.score,
                r.completed_at,
                CASE WHEN r.score >= ? THEN 1 ELSE 0 END as passed
            FROM Results r
            WHERE r.assignment_id = ?
            ORDER BY r.completed_at DESC
        """, (session_dict['qualification_marks'], session_dict['assignment_id']))
        
        attempt_columns = [column[0] for column in cursor.description]
        session_dict['attempts'] = []
        
        for attempt_row in cursor.fetchall():
            attempt_dict = dict(zip(attempt_columns, attempt_row))
            
            # Convert completed_at from UTC to IST
            if attempt_dict['completed_at']:
                attempt_dict['completed_at'] = convert_utc_to_ist(attempt_dict['completed_at'])
            
            session_dict['attempts'].append(attempt_dict)
        
        # If there are attempts, set the latest attempt info at the session level for backward compatibility
        if session_dict['attempts']:
            latest_attempt = session_dict['attempts'][0]  # First one is the latest due to DESC order
            session_dict['completed_at'] = latest_attempt['completed_at']
            session_dict['score'] = latest_attempt['score']
        
        available_sessions.append(session_dict)

    cursor.close()
    conn.close()

    return render_template('trainee_dashboard.html', sessions=available_sessions)

@app.route('/download-template/<format>')
def download_template(format):
    if format not in ['xlsx']:
        abort(404)
    return send_from_directory('static/', f'trainee_template.{format}')

@app.route('/download-excel')
def download_excel():
    if 'excel_download_path' not in session:
        flash('No excel file available for download')
        return redirect(url_for('home'))
    
    download_path = session.pop('excel_download_path')
    
    if not os.path.exists(download_path):
        flash('Excel file not found')
        return redirect(url_for('home'))
    
    try:
        return send_file(
            download_path,
            as_attachment=True,
            download_name='trainee_list_with_passwords.xlsx'
        )
    finally:
        # Clean up the file after sending
        os.unlink(download_path)

@app.route('/trainer/upload', methods=['POST'])
def upload_document():
    ALLOWED_EXTENSIONS = {'.pdf', '.docx', '.pptx', '.ppt'}
    ALLOWED_VIDEO_EXTENSIONS = {'.mp4', '.avi', '.mov'}
    ALLOWED_AUDIO_EXTENSIONS = {'.mp3'}

    try:
        qualification_marks = float(request.form.get('qualification_marks'))
        num_questions = int(request.form.get('num_questions'))
        title = request.form.get('title')
        session_id = str(uuid.uuid4())

        # Create training session first
        conn = get_db_connection()
        cursor = conn.cursor()

        try:
            # Create training session
            cursor.execute("""
                INSERT INTO TrainingSession (
                    id, title, qualification_marks, num_questions, created_by
                ) VALUES (?, ?, ?, ?, ?)
            """, (
                session_id, title, qualification_marks, num_questions, session['user_id']
            ))

            all_text_content = []
            document_sources = []
            
            # Handle document files
            if 'documents[]' in request.files:
                for file in request.files.getlist('documents[]'):
                    if file.filename:
                        document_id = str(uuid.uuid4())
                        text_content, document_url = process_document(
                            file, ALLOWED_EXTENSIONS, document_id)
                        if text_content:
                            all_text_content.append(text_content)
                            document_sources.append({
                                'id': document_id,
                                'url': document_url
                            })
                            save_training_doc(cursor, session_id, document_url, document_id)

            # Handle video files
            if 'videos[]' in request.files:
                for file in request.files.getlist('videos[]'):
                    if file.filename:
                        document_id = str(uuid.uuid4())
                        text_content, document_url = process_video(
                            file, ALLOWED_VIDEO_EXTENSIONS, document_id)
                        if text_content:
                            all_text_content.append(text_content)
                            document_sources.append({
                                'id': document_id,
                                'url': document_url
                            })
                            save_training_doc(cursor, session_id, document_url, document_id)

            # Handle audio files
            if 'audios[]' in request.files:
                for file in request.files.getlist('audios[]'):
                    if file.filename:
                        document_id = str(uuid.uuid4())
                        text_content, document_url = process_audio(
                            file, ALLOWED_AUDIO_EXTENSIONS, document_id)
                        if text_content:
                            all_text_content.append(text_content)
                            document_sources.append({
                                'id': document_id,
                                'url': document_url
                            })
                            save_training_doc(cursor, session_id, document_url, document_id)

            # Handle web URLs
            if 'web_urls[]' in request.form:
                for url in request.form.getlist('web_urls[]'):
                    if url.strip():
                        document_id = str(uuid.uuid4())
                        text_content = extract_text_from_web(url.strip())
                        if text_content:
                            all_text_content.append(text_content)
                            document_sources.append({
                                'id': document_id,
                                'url': url.strip()
                            })
                            save_training_doc(cursor, session_id, url.strip(), document_id)

            if not all_text_content:
                raise Exception('No content could be extracted from any source')

            # Combine all text content and index
            combined_text = ' '.join(all_text_content)
            chunks = chunk_text(combined_text)
            
            # Index chunks with minimal schema
            documents_to_index = [
                {
                    "id": f"{session_id}_chunk_{chunk_number}",
                    "content": chunk_content,
                    "document_url": document_sources[0]['url'],  # Using first source as reference
                    "parent_id": session_id,
                    "chunk_number": chunk_number,
                    "uploaded_at": datetime.now().isoformat() + 'Z'
                }
                for chunk_number, chunk_content in enumerate(chunks)
            ]

            # Upload in batches
            batch_size = 100
            for i in range(0, len(documents_to_index), batch_size):
                batch = documents_to_index[i:i + batch_size]
                try:
                    search_client.upload_documents(batch)
                except Exception as e:
                    print(f"Error uploading batch {i // batch_size}: {str(e)}")
                    raise

            conn.commit()
            return jsonify({
                'status': 'success',
                'session_id': session_id,
                'message': 'Training session created successfully'
            })

        except Exception as e:
            conn.rollback()
            raise Exception(f'Error creating training session: {str(e)}')
        finally:
            cursor.close()
            conn.close()

    except Exception as e:
        return jsonify({
            'status': 'error',
            'message': str(e)
        })



def save_training_doc(cursor, session_id, document_url, document_id):
    cursor.execute("""
        INSERT INTO Training_Docs (
            session_id, document_url, document_id
        ) VALUES (?, ?, ?)
    """, (session_id, document_url, document_id))

def process_document(file, allowed_extensions, document_id):
    file_extension = os.path.splitext(file.filename)[1].lower()
    if file_extension not in allowed_extensions:
        raise Exception(f'Unsupported file type. Allowed types are: {", ".join(allowed_extensions)}')
    
    temp_file = None
    try:
        # Create temporary file
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=file_extension)
        file.save(temp_file.name)
        temp_file.close()  # Close the file handle explicitly

        # Extract text based on file type
        if file_extension == '.pdf':
            text_content = extract_text_from_pdf(temp_file.name)
        elif file_extension == '.docx':
            text_content = extract_text_from_docx(temp_file.name)
        elif file_extension in ['.pptx', '.ppt']:
            text_content = extract_text_from_pptx(temp_file.name)

        # Upload to blob storage
        blob_name = f"{datetime.now().strftime('%Y%m%d_%H%M%S')}_{file.filename}"
        blob_client = blob_service_client.get_blob_client(container="media-files", blob=blob_name)
        
        with open(temp_file.name, "rb") as data:
            blob_client.upload_blob(data)
            return text_content, blob_client.url

    finally:
        # Clean up temporary file
        if temp_file and os.path.exists(temp_file.name):
            try:
                os.close(os.open(temp_file.name, os.O_RDONLY))  # Close any remaining file handles
                os.unlink(temp_file.name)
            except Exception as e:
                print(f"Warning: Failed to remove temporary file {temp_file.name}: {str(e)}")

def process_video(file, allowed_extensions, document_id):
    file_extension = os.path.splitext(file.filename)[1].lower()
    if file_extension not in allowed_extensions:
        raise Exception(f'Unsupported file type. Allowed types are: {", ".join(allowed_extensions)}')
    
    subscription_key = "4cMjRC08if3r8r2SRvuYEDVBroLmZGsiLGR5QgWPkD3VAj8PJRk6JQQJ99BBACHYHv6XJ3w3AAAAACOGQGYn"
    
    video_temp = None
    audio_temp = None
    try:
        # Create temporary file for video
        video_temp = tempfile.NamedTemporaryFile(delete=False, suffix=file_extension)
        file.save(video_temp.name)
        video_temp.close()  # Close the file handle explicitly

        # Create temporary file for audio
        audio_temp = tempfile.NamedTemporaryFile(delete=False, suffix='.mp3')
        audio_temp.close()  # Close the file handle explicitly

        # Extract audio
        video = VideoFileClip(video_temp.name)
        video.audio.write_audiofile(audio_temp.name)
        video.close()  # Close video file explicitly
        
        # Transcribe audio
        text_content = transcribe_audio(audio_temp.name, subscription_key)
        print("Uploading video to blob storage")
        blob_name = f"{datetime.now().strftime('%Y%m%d_%H%M%S')}_{file.filename}"
        blob_client = blob_service_client.get_blob_client(container="media-files", blob=blob_name)
                    
        block_list = []
        chunk_size = 4 * 1024 * 1024  # 4MB chunks

        with open(video_temp.name, "rb") as data:
            index = 0
            while True:
                chunk = data.read(chunk_size)
                if not chunk:
                    break
                block_id = base64.b64encode(str(index).zfill(6).encode()).decode()  # Base64 encode block id
                blob_client.stage_block(block_id=block_id, data=chunk)
                block_list.append({"id": block_id})
                index += 1

        try:
            blob_client.commit_block_list([block['id'] for block in block_list])  # Use only block IDs
            document_url = blob_client.url
            print("Video upload complete:", document_url)
        except Exception as e:
            print(f"Failed to commit block list: {e}")

        return text_content, document_url

    finally:
        # Clean up temporary files
        for temp_file in [video_temp, audio_temp]:
            if temp_file and os.path.exists(temp_file.name):
                try:
                    os.close(os.open(temp_file.name, os.O_RDONLY))  # Close any remaining file handles
                    os.unlink(temp_file.name)
                except Exception as e:
                    print(f"Warning: Failed to remove temporary file {temp_file.name}: {str(e)}")

def process_audio(file, allowed_extensions, document_id):
    file_extension = os.path.splitext(file.filename)[1].lower()
    if file_extension not in allowed_extensions:
        raise Exception(f'Unsupported file type. Allowed types are: {", ".join(allowed_extensions)}')
    
    subscription_key = "4cMjRC08if3r8r2SRvuYEDVBroLmZGsiLGR5QgWPkD3VAj8PJRk6JQQJ99BBACHYHv6XJ3w3AAAAACOGQGYn"
    
    temp_file = None
    try:
        # Create temporary file
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=file_extension)
        file.save(temp_file.name)
        temp_file.close()  # Close the file handle explicitly

        # Transcribe audio
        text_content = transcribe_audio(temp_file.name, subscription_key)

        print("Uploading audio to blob storage")
        blob_name = f"{datetime.now().strftime('%Y%m%d_%H%M%S')}_{file.filename}"
        blob_client = blob_service_client.get_blob_client(container="media-files", blob=blob_name)
                    
        block_list = []
        chunk_size = 4 * 1024 * 1024  # 4MB chunks

        with open(temp_file.name, "rb") as data:
            index = 0
            while True:
                chunk = data.read(chunk_size)
                if not chunk:
                    break
                block_id = base64.b64encode(str(index).zfill(6).encode()).decode()  # Base64 encode block id
                blob_client.stage_block(block_id=block_id, data=chunk)
                block_list.append({"id": block_id})
                index += 1

        try:
            blob_client.commit_block_list([block['id'] for block in block_list])  # Use only block IDs
            document_url = blob_client.url
            print("Audio upload complete:", document_url)
        except Exception as e:
            print(f"Failed to commit block list: {e}")

        return text_content, document_url

    finally:
        # Clean up temporary file
        if temp_file and os.path.exists(temp_file.name):
            try:
                os.close(os.open(temp_file.name, os.O_RDONLY))  # Close any remaining file handles
                os.unlink(temp_file.name)
            except Exception as e:
                print(f"Warning: Failed to remove temporary file {temp_file.name}: {str(e)}")

            
    
@app.route('/trainer/assign-trainees', methods=['POST'])
def assign_trainees():
    try:
        data = request.get_json()
        session_id = data.get('session_id')
        trainee_emails = data.get('trainee_emails', [])

        if not session_id or not trainee_emails:
            return jsonify({
                'status': 'error',
                'message': 'Session ID and trainee emails are required'
            })

        conn = get_db_connection()
        cursor = conn.cursor()

        try:
            # Get trainee IDs from emails
            trainee_ids = []
            for email in trainee_emails:
                cursor.execute("SELECT id FROM Users WHERE email = ?", (email.strip(),))
                result = cursor.fetchone()
                if result:
                    trainee_ids.append(result[0])

            if not trainee_ids:
                return jsonify({
                    'status': 'error',
                    'message': 'No valid trainee emails found'
                })

            # Bulk insert trainee assignments
            assignment_data = [
                (str(uuid.uuid4()), session_id, trainee_id)
                for trainee_id in trainee_ids
            ]
            
            cursor.executemany("""
                INSERT INTO TraineeAssignment (id, session_id, trainee_id)
                VALUES (?, ?, ?)
            """, assignment_data)

            conn.commit()
            return jsonify({
                'status': 'success',
                'message': 'Trainees assigned successfully'
            })

        except Exception as e:
            conn.rollback()
            raise Exception(f'Error assigning trainees: {str(e)}')
        finally:
            cursor.close()
            conn.close()

    except Exception as e:
        return jsonify({
            'status': 'error',
            'message': str(e)
        })

@app.route('/trainer/assign-trainees-file', methods=['POST'])
def assign_trainees_file():
    try:
        session_id = request.form.get('session_id')
        if 'trainee_list' not in request.files:
            return jsonify({
                'status': 'error',
                'message': 'No file uploaded'
            })

        file = request.files['trainee_list']
        if file.filename == '':
            return jsonify({
                'status': 'error',
                'message': 'No file selected'
            })

        # Process Excel file
        df = pd.read_excel(file)
        
        if 'Email' not in df.columns:
            return jsonify({
                'status': 'error',
                'message': 'File must contain an "Email" column'
            })

        trainee_emails = df['Email'].dropna().tolist()

        if not trainee_emails:
            return jsonify({
                'status': 'error',
                'message': 'No valid email addresses found in file'
            })

        # Use the existing assign_trainees logic
        return assign_trainees_with_emails(session_id, trainee_emails)

    except Exception as e:
        return jsonify({
            'status': 'error',
            'message': str(e)
        })

def assign_trainees_with_emails(session_id, trainee_emails):
    if not trainee_emails:
        return jsonify({
            'status': 'error',
            'message': 'No trainee emails provided'
        })

    # Clean and deduplicate emails
    clean_emails = list(set(email.strip().lower() for email in trainee_emails if email.strip()))
    
    if not clean_emails:
        return jsonify({
            'status': 'error',
            'message': 'No valid trainee emails found'
        })

    conn = get_db_connection()
    try:
        cursor = conn.cursor()
        cursor.fast_executemany = True

        # Create a temporary table for emails
        cursor.execute("""
        CREATE TABLE #EmailList (
            email nvarchar(255)
        )
        """)

        # Insert emails into temp table
        email_values = [(email,) for email in clean_emails]
        cursor.executemany(
            "INSERT INTO #EmailList (email) VALUES (?)",
            email_values
        )

        # Get all matching trainee IDs in a single query
        trainee_query = """
        SELECT id 
        FROM Users u
        INNER JOIN #EmailList e ON LOWER(u.email) = e.email
        """
        
        cursor.execute(trainee_query)
        trainee_ids = [row[0] for row in cursor.fetchall()]

        if not trainee_ids:
            return jsonify({
                'status': 'error',
                'message': 'No valid trainee emails found in the system'
            })

        # Check for existing assignments
        check_existing = """
        SELECT trainee_id 
        FROM TraineeAssignment 
        WHERE session_id = ? 
        AND trainee_id IN (
            SELECT u.id 
            FROM Users u 
            INNER JOIN #EmailList e ON LOWER(u.email) = e.email
        )
        """
        
        cursor.execute(check_existing, (session_id,))
        existing_trainee_ids = {row[0] for row in cursor.fetchall()}

        # Filter out already assigned trainees
        new_trainee_ids = [tid for tid in trainee_ids if tid not in existing_trainee_ids]

        if not new_trainee_ids:
            return jsonify({
                'status': 'success',
                'message': 'All trainees were already assigned to this session'
            })

        # Prepare assignment data
        assignment_data = [
            (str(uuid.uuid4()), session_id, trainee_id)
            for trainee_id in new_trainee_ids
        ]

        # Bulk insert assignments
        insert_query = """
        INSERT INTO TraineeAssignment (id, session_id, trainee_id)
        VALUES (?, ?, ?)
        """
        
        # Use batched insert for better performance
        BATCH_SIZE = 1000
        for i in range(0, len(assignment_data), BATCH_SIZE):
            batch = assignment_data[i:i + BATCH_SIZE]
            cursor.executemany(insert_query, batch)
            conn.commit()

        # Clean up temp table
        cursor.execute("DROP TABLE #EmailList")
        
        return jsonify({
            'status': 'success',
            'message': f'Successfully assigned {len(new_trainee_ids)} trainees'
        })

    except Exception as e:
        conn.rollback()
        app.logger.error(f'Error in assign_trainees_with_emails: {str(e)}')
        return jsonify({
            'status': 'error',
            'message': 'Error assigning trainees. Please try again.'
        })
    finally:
        # Ensure temp table is dropped even if there's an error
        try:
            cursor.execute("IF OBJECT_ID('tempdb..#EmailList') IS NOT NULL DROP TABLE #EmailList")
        except:
            pass
        conn.close()
    

@app.route('/trainer/view_attempt/<assignment_id>/<result_id>')
def view_attempt(assignment_id, result_id):
    if 'user_id' not in session:
        return redirect(url_for('login'))

    conn = get_db_connection()
    cursor = conn.cursor()

    # Get the assignment details
    cursor.execute("""
        SELECT 
            ta.trainee_id,
            u.name as trainee_name,
            u.email as trainee_email,
            ts.qualification_marks,
            ts.num_questions
        FROM TraineeAssignment ta
        JOIN TrainingSession ts ON ta.session_id = ts.id
        JOIN Users u ON ta.trainee_id = u.id
        WHERE ta.id = ?
    """, (assignment_id,))

    assignment = cursor.fetchone()
    if not assignment:
        cursor.close()
        conn.close()
        flash('No assignment found for this attempt.')
        return redirect(url_for('trainer_dashboard'))

    # Get the questions and responses for the specific result
    cursor.execute("""
        SELECT 
            tq.question_text,
            tq.option_a,
            tq.option_b,
            tq.option_c,
            tq.option_d,
            tq.response,
            tq.correct_answer
        FROM TestQuestions tq
        WHERE tq.test_id = (
            SELECT test_id FROM Results WHERE id = ?
        )
        ORDER BY tq.id
    """, (result_id,))

    questions = cursor.fetchall()

    cursor.close()
    conn.close()

    return render_template('view_attempt.html', assignment=assignment, questions=questions)


@app.route('/view_results/<session_id>')
def view_results(session_id):
    if 'user_id' not in session:
        return redirect(url_for('login'))

    conn = get_db_connection()
    cursor = conn.cursor()

    # Get trainee assignment and session details
    cursor.execute("""
        SELECT 
            ta.id as assignment_id,
            ta.trainee_id,
            ts.qualification_marks,
            ts.num_questions,
            u.name as trainee_name,
            u.email as trainee_email
        FROM TraineeAssignment ta
        JOIN TrainingSession ts ON ta.session_id = ts.id
        JOIN Users u ON ta.trainee_id = u.id
        WHERE ts.id = ? AND ta.trainee_id = ?
    """, (session_id, session['user_id']))

    assignment = cursor.fetchone()
    if not assignment:
        cursor.close()
        conn.close()
        flash('No assignment found for this test session.')
        return redirect(url_for('trainee_dashboard'))

    # Get the latest test result
    cursor.execute("""
        SELECT TOP 1
            r.id as result_id,
            r.test_id,
            r.score,
            r.completed_at,
            r.completed
        FROM Results r
        WHERE r.assignment_id = ?
        ORDER BY r.completed_at DESC
    """, (assignment.assignment_id,))

    result = cursor.fetchone()

    # Get questions and answers for the latest test attempt
    test_details = None
    if result:
        cursor.execute("""
            SELECT 
                id,
                question_text,
                option_a,
                option_b,
                option_c,
                option_d,
                response,
                correct_answer
            FROM TestQuestions
            WHERE test_id = ?
            ORDER BY id
        """, (result.test_id,))

        questions = cursor.fetchall()

        test_details = {
            'result': result,
            'questions': questions,
            'num_correct': sum(1 for q in questions if q.response == q.correct_answer),
            'total_questions': len(questions),
            'passed': result.score >= assignment.qualification_marks if result.score else False
        }

    # Get all previous attempts
    cursor.execute("""
        SELECT 
            r.id,
            r.score,
            r.completed_at,
            r.completed
        FROM Results r
        WHERE r.assignment_id = ?
        ORDER BY r.completed_at DESC
    """, (assignment.assignment_id,))

    previous_attempts = cursor.fetchall()

    cursor.close()
    conn.close()

    return render_template(
        'results.html',
        assignment=assignment,
        test_details=test_details,
        previous_attempts=previous_attempts,
        qualification_marks=assignment.qualification_marks
    )

@app.route('/trainee/test/<session_id>')
def take_test(session_id):
    conn = get_db_connection()
    cursor = conn.cursor()
    
    # Get training session and assignment info
    cursor.execute("""
        SELECT 
            ts.id,
            ts.title,
            ts.qualification_marks,
            ts.num_questions,
            ta.id as assignment_id
        FROM TrainingSession ts
        JOIN TraineeAssignment ta ON ts.id = ta.session_id
        WHERE ts.id = ? AND ta.trainee_id = ?
    """, (session_id, session['user_id']))
    
    result = cursor.fetchone()
    print(result)
    if not result:
        cursor.close()
        conn.close()
        return redirect(url_for('trainee_dashboard'))

    # Search for document chunks using session_id as parent_id
    
    filter_condition = f"parent_id eq '{session_id.lower()}'"
    results = search_client.search(
        search_text="*",
        filter=filter_condition,
        select="content,chunk_number",
        order_by="chunk_number"
    )

    full_content = ""
    try:
        documents = list(results)
        if not documents:
            cursor.close()
            conn.close()
            flash('Document content not found')
            return redirect(url_for('trainee_dashboard'))

        full_content = " ".join(doc['content'] for doc in documents)
        print(full_content)

    except Exception as e:
        cursor.close()
        conn.close()
        print(f"Error retrieving document chunks: {str(e)}")
        flash('Error retrieving document content')
        return redirect(url_for('trainee_dashboard'))
    
    # Generate a new test_id for this attempt
    test_id = str(uuid.uuid4())

    # Generate questions and store them with the new test_id
    questions = generate_questions(full_content, result.num_questions, test_id)

    # Store test information in session
    session['current_test'] = {
        'session_id': session_id,
        'assignment_id': result.assignment_id,
        'test_id': test_id,
        'questions': questions
    }

    cursor.close()
    conn.close()

    return render_template(
        'test.html',
        questions=questions,
        session_id=session_id,
        num_questions=result.num_questions,
        title=result.title
    )

@app.route('/trainee/submit_test', methods=['POST'])
def submit_test():
    session_id = request.form.get('session_id')
    current_test = session.get('current_test')

    if not current_test or current_test['session_id'] != session_id:
        flash('Test session expired')
        return redirect(url_for('trainee_dashboard'))

    conn = get_db_connection()
    cursor = conn.cursor()

    assignment_id = current_test['assignment_id']
    test_id = current_test['test_id']
    answers = []
    correct_count = 0
    total_questions = len(current_test['questions'])

    # Process each answer and update TestQuestions
    for i in range(total_questions):
        answer = request.form.get(f'question_{i}')
        answers.append(answer)

        if answer == current_test['questions'][i]['correct_answer']:
            correct_count += 1

        cursor.execute("""
            UPDATE TestQuestions 
            SET response = ?
            WHERE test_id = ? AND question_text = ?
        """, (
            answer,
            test_id,
            current_test['questions'][i]['text']
        ))

    score = (correct_count / total_questions) * 100

    # Create new result record
    result_id = str(uuid.uuid4())
    cursor.execute("""
        INSERT INTO Results (
            id, assignment_id, test_id, completed, score, completed_at
        ) VALUES (?, ?, ?, 1, ?, GETDATE())
    """, (result_id, assignment_id, test_id, score))

    conn.commit()
    cursor.close()
    conn.close()

    session['user_answers'] = answers
    session.pop('current_test', None)

    return redirect(url_for('view_results', session_id=session_id))


@app.route('/trainer/download_trainee_details/<session_id>')
def download_trainee_details(session_id):
    if 'user_id' not in session:
        return redirect(url_for('login'))
    
    conn = get_db_connection()
    cursor = conn.cursor()
    
    # Modified query to get ALL trainees assigned to the session
    cursor.execute("""
        SELECT 
            u.name,
            u.email,
            COALESCE(MAX(r.score), 0) as best_score,
            COUNT(r.id) as attempts,
            CASE 
                WHEN MAX(r.score) >= ts.qualification_marks THEN 'Passed'
                WHEN COUNT(r.id) > 0 THEN 'Failed'
                ELSE 'Not Attempted'
            END as status,
            MAX(r.completed_at) as last_attempt,
            ts.qualification_marks,
            ts.created_at
        FROM TraineeAssignment ta
        JOIN Users u ON ta.trainee_id = u.id
        JOIN TrainingSession ts ON ta.session_id = ts.id
        LEFT JOIN Results r ON ta.id = r.assignment_id
        WHERE ta.session_id = ?
        GROUP BY u.name, u.email, ts.qualification_marks, ts.created_at
    """, (session_id,))
    
    # Convert rows to list of dictionaries
    columns = ['Trainee Name', 'Email', 'Best Score (%)', 'Number of Attempts', 
               'Status', 'Last Attempt', 'Qualification Marks (%)', 'Session Created']
    
    rows_data = []
    for row in cursor.fetchall():
        row_list = list(row)
        # Format the score to 2 decimal places
        row_list[2] = round(float(row_list[2]), 2) if row_list[2] is not None else 0.0
        rows_data.append(dict(zip(columns, row_list)))

    # Create DataFrame
    df = pd.DataFrame(rows_data)
    
    # Format dates
    df['Last Attempt'] = pd.to_datetime(df['Last Attempt']).apply(
        lambda x: convert_utc_to_ist(x).strftime('%Y-%m-%d %H:%M') if pd.notnull(x) else 'Not Attempted'
    )
    df['Session Created'] = pd.to_datetime(df['Session Created']).apply(
        lambda x: convert_utc_to_ist(x).strftime('%Y-%m-%d %H:%M')
    )
    
    # Create Excel file in memory
    output = BytesIO()
    with pd.ExcelWriter(output, engine='xlsxwriter') as writer:
        df.to_excel(writer, sheet_name='Trainee Details', index=False)
        
        # Get workbook and worksheet objects
        workbook = writer.book
        worksheet = writer.sheets['Trainee Details']
        
        # Add formats
        header_format = workbook.add_format({
            'bold': True,
            'bg_color': '#4B88CB',
            'font_color': 'white',
            'border': 1
        })
        
        # Add conditional formatting for status column
        status_col = columns.index('Status')
        worksheet.conditional_format(1, status_col, len(df) + 1, status_col, {
            'type': 'cell',
            'criteria': 'equal to',
            'value': '"Passed"',
            'format': workbook.add_format({'bg_color': '#C6EFCE', 'font_color': '#006100'})
        })
        worksheet.conditional_format(1, status_col, len(df) + 1, status_col, {
            'type': 'cell',
            'criteria': 'equal to',
            'value': '"Failed"',
            'format': workbook.add_format({'bg_color': '#FFC7CE', 'font_color': '#9C0006'})
        })
        worksheet.conditional_format(1, status_col, len(df) + 1, status_col, {
            'type': 'cell',
            'criteria': 'equal to',
            'value': '"Not Attempted"',
            'format': workbook.add_format({'bg_color': '#FFEB9C', 'font_color': '#9C6500'})
        })
        
        # Format the header
        for col_num, value in enumerate(df.columns.values):
            worksheet.write(0, col_num, value, header_format)
            worksheet.set_column(col_num, col_num, len(value) + 5)
    
    output.seek(0)
    cursor.close()
    conn.close()
    
    return send_file(
        output,
        mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
        as_attachment=True,
        download_name=f'trainee_details_{session_id}.xlsx'
    )

# Modified main entry point for Azure
if __name__ == '__main__':
    app.run(host='0.0.0.0', port=int(os.environ.get('PORT', 5000)))
else:
    application = app
